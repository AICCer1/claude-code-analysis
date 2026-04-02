from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

OUTDIR = Path('/root/.openclaw/workspace/claude-code-analysis/ppt')
OUTDIR.mkdir(parents=True, exist_ok=True)

BLUE = RGBColor(32, 84, 147)
LIGHT_BLUE = RGBColor(225, 236, 248)
DARK = RGBColor(33, 37, 41)
GRAY = RGBColor(107, 114, 128)
GREEN = RGBColor(30, 120, 70)
LIGHT_GREEN = RGBColor(230, 244, 234)
ORANGE = RGBColor(196, 110, 0)
LIGHT_ORANGE = RGBColor(252, 239, 213)
RED = RGBColor(160, 48, 48)
LIGHT_RED = RGBColor(251, 232, 232)
WHITE = RGBColor(255,255,255)


def add_title(slide, title, subtitle=None):
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.8), Inches(0.7))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = DARK
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.62), Inches(0.95), Inches(11), Inches(0.35))
        p2 = tx2.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = GRAY


def add_footer(slide, text='外部 AI 资源辅助内部研发流程'):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(7.05), Inches(12.1), Inches(0.02))
    line.fill.solid(); line.fill.fore_color.rgb = LIGHT_BLUE; line.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(5), Inches(0.2))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text; r.font.size = Pt(8.5); r.font.color.rgb = GRAY


def add_bullets_box(slide, x, y, w, h, title, bullets, fill=LIGHT_BLUE, title_color=BLUE):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = fill
    tb = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.12), w - Inches(0.25), h - Inches(0.18))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = title_color
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.space_before = Pt(3)
    return box


def add_kpi_box(slide, x, y, w, h, title, value, note, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.color.rgb = fill
    tb = slide.shapes.add_textbox(x + Inches(0.16), y + Inches(0.14), w - Inches(0.2), h - Inches(0.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title; r.font.size = Pt(11); r.font.color.rgb = GRAY; r.font.bold = True
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = value; r2.font.size = Pt(20); r2.font.bold = True; r2.font.color.rgb = DARK
    p3 = tf.add_paragraph(); r3 = p3.add_run(); r3.text = note; r3.font.size = Pt(9.5); r3.font.color.rgb = DARK


def add_step(slide, x, y, w, h, title, body, fill, line_color=None, title_color=WHITE, body_color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color or fill
    tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.08), w - Inches(0.2), h - Inches(0.16))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = title_color
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = body; r2.font.size = Pt(9.5); r2.font.color.rgb = body_color
    return shp


def connect(slide, x1, y1, x2, y2, color=GRAY, width=1.75):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def make_report_deck(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, '外部 AI 资源辅助内部研发', '两页版 / 汇报风：强调价值、流程与信息安全边界')
    add_bullets_box(slide, Inches(0.6), Inches(1.35), Inches(4.05), Inches(3.15),
                    '外网阶段', [
                        '使用 Codex、Claude Code 进行新技术 / 新功能 / 新特性的 demo 级快速验证',
                        '使用 OpenClaw + Exa 检索最新论文、产品文档、开源仓库与实现案例',
                        '自动完成代码仓创建、推送、外部仓库分析拆解、模块提炼与方案沉淀',
                        '输出实现细节、迁移方案与参考代码仓'
                    ], fill=LIGHT_BLUE, title_color=BLUE)
    add_bullets_box(slide, Inches(4.85), Inches(1.35), Inches(4.05), Inches(3.15),
                    '内网阶段', [
                        '下载外部阶段产出的代码仓与资料',
                        '由内网模型继续开发与代码生成',
                        '结合内部项目进行集成、联调与落地',
                        '形成内部可维护、可持续演进的实现'
                    ], fill=LIGHT_GREEN, title_color=GREEN)
    add_bullets_box(slide, Inches(9.1), Inches(1.35), Inches(3.6), Inches(3.15),
                    '核心价值', [
                        '缩短新技术验证周期',
                        '降低前期调研和方案探索成本',
                        '把外部探索转化为内部研发加速器',
                        '保证数据单向流动，兼顾效率与安全'
                    ], fill=LIGHT_ORANGE, title_color=ORANGE)

    add_kpi_box(slide, Inches(0.7), Inches(5.0), Inches(2.7), Inches(1.2), '定位', '外探内用', '外网负责探索验证，内网负责正式开发与集成', LIGHT_BLUE)
    add_kpi_box(slide, Inches(3.7), Inches(5.0), Inches(2.7), Inches(1.2), '数据方向', '外 → 内', '信息只从外部进入内部，避免敏感数据外流', LIGHT_GREEN)
    add_kpi_box(slide, Inches(6.7), Inches(5.0), Inches(2.7), Inches(1.2), '交付物', '方案 + 代码仓', '不仅有调研结论，也有可继续开发的代码基础', LIGHT_ORANGE)
    add_kpi_box(slide, Inches(9.7), Inches(5.0), Inches(2.3), Inches(1.2), '目标', '研发加速', '服务内部研发团队快速落地', LIGHT_RED)
    add_footer(slide)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, '流程图与信息安全边界', '汇报风：流程清晰、边界明确、适合直接讲解')
    steps_y = Inches(2.0)
    xs = [Inches(0.55), Inches(2.85), Inches(5.15), Inches(7.45), Inches(9.75)]
    w = Inches(2.0); h = Inches(1.2)
    add_step(slide, xs[0], steps_y, w, h, '需求输入', '新技术\n新功能\n新特性验证', BLUE)
    add_step(slide, xs[1], steps_y, w, h, '外网探索', 'Codex / Claude Code\nOpenClaw + Exa', ORANGE)
    add_step(slide, xs[2], steps_y, w, h, '形成产出', '论文 / 文档 / 仓库分析\nDemo / 方案 / 参考代码', GREEN)
    add_step(slide, xs[3], steps_y, w, h, '内网承接', '下载代码仓与方案\n内部模型继续开发', BLUE)
    add_step(slide, xs[4], steps_y, w, h, '项目集成', '接入内部项目\n联调 / 测试 / 落地', GREEN)
    for i in range(4):
        connect(slide, xs[i] + w, steps_y + h/2, xs[i+1], steps_y + h/2, GRAY)

    sec = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.25), Inches(12.0), Inches(1.65))
    sec.fill.solid(); sec.fill.fore_color.rgb = LIGHT_RED; sec.line.color.rgb = LIGHT_RED
    tb = slide.shapes.add_textbox(Inches(0.95), Inches(4.45), Inches(11.5), Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]; r = p.add_run(); r.text = '信息安全边界'; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = RED
    bullets = [
        '外网环境：用于技术洞察、公开资料检索、代码仓分析与 demo 级验证，不引入内部敏感数据',
        '内网环境：用于正式开发、与内部项目集成、代码沉淀与交付',
        '数据方向保持为“外 → 内”，避免内部信息向外部模型或外部环境回流'
    ]
    for b in bullets:
        p = tf.add_paragraph(); p.text = b; p.font.size = Pt(11); p.font.color.rgb = DARK
    add_footer(slide)
    prs.save(path)


def make_method_deck(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, '外部 AI 资源辅助内部研发：方法拆解', '两页版 / 方法论风：强调工具链、职责分层与交付物')

    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.45), Inches(5.85), Inches(4.95))
    left.fill.solid(); left.fill.fore_color.rgb = LIGHT_BLUE; left.line.color.rgb = LIGHT_BLUE
    right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(1.45), Inches(5.85), Inches(4.95))
    right.fill.solid(); right.fill.fore_color.rgb = LIGHT_GREEN; right.line.color.rgb = LIGHT_GREEN

    tb1 = slide.shapes.add_textbox(Inches(0.95), Inches(1.68), Inches(5.3), Inches(4.4))
    tf1 = tb1.text_frame
    p = tf1.paragraphs[0]; r = p.add_run(); r.text = '外网职责'; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = BLUE
    for b in [
        '工具：Codex、Claude Code、OpenClaw + Exa',
        '任务：最新技术洞察、论文/文档检索、竞品/开源仓分析、demo 级快速实现',
        '产出：实现思路、迁移方案、参考代码、外部代码仓',
        '特点：更适合做探索、验证、拆解、生成初版方案'
    ]:
        p = tf1.add_paragraph(); p.text = b; p.font.size = Pt(11.5); p.font.color.rgb = DARK

    tb2 = slide.shapes.add_textbox(Inches(7.0), Inches(1.68), Inches(5.3), Inches(4.4))
    tf2 = tb2.text_frame
    p = tf2.paragraphs[0]; r = p.add_run(); r.text = '内网职责'; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = GREEN
    for b in [
        '工具：公司内网模型与内部开发环境',
        '任务：下载外部产物、继续实现、与内部项目集成、联调与沉淀',
        '产出：内部可维护代码、与业务系统适配的实现、研发资产沉淀',
        '特点：更适合做正式开发、项目集成与安全可控的交付'
    ]:
        p = tf2.add_paragraph(); p.text = b; p.font.size = Pt(11.5); p.font.color.rgb = DARK

    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.55), Inches(6.25), Inches(4.25), Inches(0.52))
    tag.fill.solid(); tag.fill.fore_color.rgb = WHITE; tag.line.color.rgb = GRAY
    tb = slide.shapes.add_textbox(Inches(4.75), Inches(6.33), Inches(3.9), Inches(0.2))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = '方法核心：外部探索 + 内部落地 + 单向数据流'; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = DARK
    add_footer(slide)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, '端到端流程与交付物', '方法论风：强调每一阶段的输入、动作与产出')

    y = Inches(1.55)
    w = Inches(2.15); h = Inches(1.15)
    blocks = [
        ('输入', '研发问题\n技术目标\n验证需求', LIGHT_BLUE, BLUE),
        ('探索', '检索论文 / 文档\n分析外部仓库\n形成理解', LIGHT_ORANGE, ORANGE),
        ('验证', '快速 demo\n功能原型\n特性试验', LIGHT_GREEN, GREEN),
        ('沉淀', '方案文档\n迁移说明\nGitHub 仓库', LIGHT_BLUE, BLUE),
        ('承接', '内网下载\n内部模型继续开发\n项目集成', LIGHT_GREEN, GREEN),
    ]
    xpos = [Inches(0.55), Inches(2.95), Inches(5.35), Inches(7.75), Inches(10.15)]
    for (title, body, fill, color), x in zip(blocks, xpos):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.color.rgb = color
        tb = slide.shapes.add_textbox(x + Inches(0.08), y + Inches(0.1), w - Inches(0.16), h - Inches(0.15))
        tf = tb.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = color
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = body; r2.font.size = Pt(10.5); r2.font.color.rgb = DARK
    for i in range(4):
        connect(slide, xpos[i] + w, y + h/2, xpos[i+1], y + h/2, GRAY)

    # lower lane
    lane = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(4.05), Inches(12.0), Inches(2.0))
    lane.fill.solid(); lane.fill.fore_color.rgb = WHITE; lane.line.color.rgb = LIGHT_BLUE
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(4.25), Inches(11.5), Inches(1.55))
    tf = tb.text_frame
    p = tf.paragraphs[0]; r = p.add_run(); r.text = '交付物与边界'; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = DARK
    for b in [
        '外网输出的交付物可以是：技术洞察、方案报告、参考实现、可继续开发的代码仓。',
        '内网承接的是外部公开信息和产出物，不传递内部敏感数据到外网环境。',
        '整个方法的关键不是“直接把外部代码搬进内网”，而是“外部高效探索、内部安全落地”。'
    ]:
        p = tf.add_paragraph(); p.text = b; p.font.size = Pt(11); p.font.color.rgb = DARK
    add_footer(slide)
    prs.save(path)


make_report_deck(OUTDIR / '外部AI资源辅助内部研发_汇报风_2页.pptx')
make_method_deck(OUTDIR / '外部AI资源辅助内部研发_方法论风_2页.pptx')
print('generated')
